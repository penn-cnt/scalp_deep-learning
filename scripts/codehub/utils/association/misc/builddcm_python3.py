#!/usr/bin/env python
import sys,os
import argparse
import numpy as np
from pptx import Presentation

zeropadjoin=lambda a_b:'%s%02d'%(a_b[0],int(a_b[1])) if all([a_b[0],a_b[1]]) else None
def build_lines(myid,electrodes,gnd,ref):
    header_lines=['# %s channel mapping'%myid,'']
    fname='%s_channelMapping.txt'%myid

    #build and sort the lines
    idx_gnd,idx_ref,last,empties,elec_lines=None,None,None,[],[]
    for idx,name in enumerate(electrodes):
        if name=='':empties.append(idx)
        elif gnd in [name]:idx_gnd=idx
        elif ref in [name]:idx_ref=idx
        else:
            last=idx
            elec_lines.append('%d %s'%(idx,name))

    empties=[idx for idx in empties if idx<last] # remove empties after last elec
    if empties:header_lines.extend(['# empty %s'%' '.join([str(empty) for empty in empties]),''])
    if None==gnd:gnd='?'
    if None==ref:ref='?'
    # insert ground and ref
    helper=lambda a_b_c:'# %d %s %s'%(a_b_c[0],a_b_c[1],a_b_c[2]) if a_b_c[0] else '# %s %s'%(a_b_c[1],a_b_c[2])
    header_lines.extend(list(map(helper,[(idx_gnd,gnd,'GND'),(idx_ref,ref,'REF')]))+[''])
    return fname,header_lines+elec_lines

def split_elec(mystr):
    a,b=[],[]
    for ind,ch in enumerate(mystr):
        if ch.isalpha():a.append(ch)
        else:break
    if len(mystr)==len(a):return ''.join(a).upper(),''
    b=mystr[ind:]
    return ''.join(a).upper(),b.strip()

def parse_elec(mystr):
    if ''==mystr: return ['']
    elec,numstr=split_elec(mystr)
    if not numstr: return [elec]
    if not any([x in numstr for x in [',','-']]): return [''.join([elec,numstr])]
    ints,words=[],[x.strip() for x in numstr.split(',')]

    for word in words:
        if word: # if not excluded
            try:
                ints.append(int(word))
            except ValueError:
                a,b=list(map(int,word.split('-')))
                ints.extend([e for e in range(a,b+1)])
        else: ints.append(word)
    if not elec in ['EKG','ECG']: return ["%s%02d" % (elec,i) if i else i for i in ints]
    else: return ["%s%d" %(elec,i) if i else i for i in ints]

def manual_entry(myid, out, gnd, ref):

    print('Please enter electrodes by name in order, one lead at a time.')
    print('\te.g. LA1-4,6-10\n\t or EKG1-2\n\t or C3')
    print('Simply hit return for any gaps.')
    print('Type ^c when done.')

    # Iterate over entries until manual break
    electrodes=[]
    while True:
        try:
            raw_val = str(input('electrodes (%d): '%(len(electrodes)+1)))
            electrodes.extend(parse_elec(raw_val))
        except KeyboardInterrupt:
            print('\n')
            break
    return electrodes

def ppt_read(ppt_file,nlead):

    # Get the user to add lead info
    if nlead == None:
        nlead = int(input("Please provide number of leads in the powerpoint file: "))

    # Read in the file provided by the user
    prs = Presentation(ppt_file)

    # Get the slides in memory
    slides_list = []
    for slide in prs.slides:
        slides_list.append(slide)

    # Calculate the number of rows that are complete and the remainder
    ncol   = 16
    nrows  = np.floor(nlead/ncol).astype('int')
    ntail  = nlead%ncol
    rowcnt = 0

    # Parse the channel pages, save to dict. (Leads are supposed to always be the last four pages)
    chmap     = {}
    electrodes = []
    for slide_cnt,slide in enumerate(slides_list[-4:]):
        
        # Read in the table data
        output = []
        for shape in slide.shapes:
            values = []
            if hasattr(shape, "table"):
                cells = shape.table.iter_cells()
                for icell in cells:
                    values.append(icell.text)

            # Clean up the ocassional character array (versus strings)
            flag = True
            for idx,ivalue in enumerate(values):
                if ivalue == '':
                    if flag:
                        output.append(ivalue)
                        flag = False
                else:
                    output.append(ivalue)
                    flag = True

        # Attempt to clean up the grnd and ref cells
        try:
            output.pop(0)
            if slide_cnt in [0,2]: output.pop(ncol)
            output.pop(2*ncol)
            output.pop(4*ncol)
            output.pop(6*ncol)
            if slide_cnt in [0,2]: output.pop(7*ncol)
        except IndexError:
            pass

        # Iterate over number of rows on this page, checking against maximum number
        while len(output) > 0:
            if rowcnt < nrows:
                channels = output[:ncol]
                output   = output[ncol:]
                leads    = output[:ncol]
                output   = output[ncol:]
                rowcnt  += 1
            else:
                output.pop(ntail)
                channels = output[:ntail]
                output   = output[ntail:]
                leads    = output[:ntail]
                output   = []
            for idx,ichannel in enumerate(channels):
                chmap[ichannel] = leads[idx]
                electrodes.append(f"{ichannel.replace(" ", "")}-{leads[idx]}")

    return electrodes

def main(myid=None, out=None, gnd=None, ref=None, ppt_file=None, nlead=None):
    
    # Obtain patient id if calld as a function. 
    if not myid:
        myid = str(input("patient ID: ")).strip()
    myid = myid.upper()
    
    # Call electrode naming functions
    if ppt_file == None:
        electrodes = manual_entry(myid, out, gnd, ref)
    else:
        electrodes = ppt_read(ppt_file,nlead)

    # Output logic
    if electrodes:
        # get ground and ref
        if not all([gnd,ref]):
            gnd,ref=[input(prompt) for prompt in ['ground? ','reference? ']]
        gnd,ref=[zeropadjoin(split_elec(e)) for e in [gnd,ref]]
        # write file
        fname,lines=build_lines(myid,electrodes,gnd,ref)
        if out != None:
            with open(os.path.join(os.path.abspath(out),fname),'w') as f:
              f.write('\n'.join(lines))
        else:
            with open(os.path.join(os.getcwd(),fname),'w') as f:
                f.write('\n'.join(lines))

if __name__=='__main__':
    
    """
    Unattributed creation.

    Minor refactorization, comments, and clean up, Brian Prager. 12/15/2023.
    """

    # Argument parser
    parser = argparse.ArgumentParser()
    parser.add_argument("id", type=str, help="id", default=None)
    parser.add_argument("--out", dest="out", type=str, help="output", default=None)
    parser.add_argument("--gnd", dest="gnd", type=str, help="ground", default=None)
    parser.add_argument("--ref", dest="ref", type=str, help="reference", default=None)
    parser.add_argument("--ppt_file", dest="ppt_file", type=str, help="Path to powerpoint file. If provided, automatically assign channels.", default=None)
    parser.add_argument("--nlead", dest="nlead", type=int, help="Number of leads in the powerpoint file. Needed for auto assigning.", default=None)
    args = parser.parse_args()

    # Call the main function
    main(myid=args.id, out=args.out, gnd=args.gnd, ref=args.ref, ppt_file=args.ppt_file, nlead=args.nlead)
