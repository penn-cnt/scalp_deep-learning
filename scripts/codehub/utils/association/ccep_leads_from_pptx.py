import numpy as np
from sys import argv,exit
from pptx import Presentation

if __name__ == '__main__':

    # Read in the file provided by the user
    prs   = Presentation(argv[1])
    nlead = int(argv[2]) 

    # Get the slides in memory
    slides_list = []
    for slide in prs.slides:
        slides_list.append(slide)

    # Calculate the number of rows that are complete and the remainder
    ncol   = 16
    nrows  = np.floor(nlead/ncol).astype('int')
    ntail  = nlead%ncol
    rowcnt = 0

    # Parse the channel pages, save to dict. (Leads are supposed to always be the last four pages)
    chmap = {}
    for slide_cnt,slide in enumerate(slides_list[-4:]):
        
        # Read in the table data
        output = []
        for shape in slide.shapes:
            values = []
            if hasattr(shape, "table"):
                cells = shape.table.iter_cells()
                for icell in cells:
                    values.append(icell.text)

            # Clean up the ocassional character array (versus strings)
            flag = True
            for idx,ivalue in enumerate(values):
                if ivalue == '':
                    if flag:
                        output.append(ivalue)
                        flag = False
                else:
                    output.append(ivalue)
                    flag = True

        # Attempt to clean up the grnd and ref cells
        try:
            output.pop(0)
            if slide_cnt in [0,2]: output.pop(ncol)
            output.pop(2*ncol)
            output.pop(4*ncol)
            output.pop(6*ncol)
            if slide_cnt in [0,2]: output.pop(7*ncol)
        except IndexError:
            pass

        # Iterate over number of rows on this page, checking against maximum number
        while len(output) > 0:
            if rowcnt < nrows:
                channels = output[:ncol]
                output   = output[ncol:]
                leads    = output[:ncol]
                output   = output[ncol:]
                rowcnt  += 1
            else:
                output.pop(ntail)
                channels = output[:ntail]
                output   = output[ntail:]
                leads    = output[:ntail]
                output   = []
            for idx,ichannel in enumerate(channels):
                chmap[ichannel] = leads[idx]

